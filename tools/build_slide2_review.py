from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DL=Path(r"C:\Users\abdul\Downloads")
SRC=DL/"KSP Datathon 2026 _ Prototype Submission Template.pptx"
OUT=DL/"NAMMA_KSP_Slide_02_Review.pptx"
NAVY=RGBColor(16,42,86); TEAL=RGBColor(11,124,134); GOLD=RGBColor(214,167,44)
INK=RGBColor(27,42,62); MUTED=RGBColor(91,107,128); WHITE=RGBColor(255,255,255)
LINE=RGBColor(214,224,235); PALE=RGBColor(244,247,250)

prs=Presentation(SRC)
s=prs.slides[1]
for sh in list(s.shapes)[1:]:
    sh._element.getparent().remove(sh._element)

def rect(x,y,w,h,fill=WHITE,line=LINE,radius=False):
    q=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                         Inches(x),Inches(y),Inches(w),Inches(h))
    q.fill.solid(); q.fill.fore_color.rgb=fill; q.line.color.rgb=line; q.line.width=Pt(.7)
    return q
def text(v,x,y,w,h,size=14,color=INK,bold=False,align=PP_ALIGN.LEFT):
    q=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); f=q.text_frame
    f.clear(); f.word_wrap=True; f.margin_left=f.margin_right=0; f.margin_top=f.margin_bottom=0
    f.vertical_anchor=MSO_ANCHOR.TOP; p=f.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=v; r.font.name="Aptos Display" if bold else "Aptos"; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return q

# Fixed title zone: nothing else is allowed above y=2.0.
text("BRIEF ABOUT THE SOLUTION",.66,.84,3.0,.22,8,GOLD,True)
text("Natural-language crime intelligence, grounded in evidence",.66,1.14,8.7,.48,23,NAVY,True)
text("One governed workspace for investigators, analysts, supervisors and policymakers.",.66,1.67,8.7,.25,10,MUTED)
rect(.66,1.96,.78,.035,TEAL,TEAL)

# Main proposition.
rect(.66,2.26,5.34,2.52,NAVY,NAVY,True)
text("THE PROPOSITION",.94,2.53,1.65,.22,8,GOLD,True)
text("Turn fragmented records into connected, explainable intelligence.",.94,2.91,4.55,.72,23,WHITE,True)
text("Ask in English or Kannada. Trace incidents, people and places. Surface patterns, risk and defensible next actions.",.94,3.92,4.52,.52,11,RGBColor(221,231,243))

# Three proof points with generous spacing and no nested cards.
proof=[("5,000","synthetic FIR records",NAVY),("7","linked datasets",TEAL),("2","languages + voice",GOLD)]
for i,(n,label,c) in enumerate(proof):
    y=2.26+i*.84
    rect(6.28,y,3.05,.65,WHITE,LINE,False)
    rect(6.28,y,.07,.65,c,c,False)
    text(n,6.55,y+.11,.75,.3,20,c,True)
    text(label,7.34,y+.19,1.72,.22,10,NAVY,True)

text("ASK  →  CONNECT  →  EXPLAIN  →  ACT",6.3,4.82,3.0,.22,9,TEAL,True,PP_ALIGN.CENTER)
prs.save(OUT)
print(OUT)
