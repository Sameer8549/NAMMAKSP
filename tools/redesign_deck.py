from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

DL=Path(r"C:\Users\abdul\Downloads")
ROOT=Path(r"C:\Users\abdul\OneDrive\Documents\NAMMAKSP")
SRC=DL/"NAMMA_KSP_World_Class_Submission_Deck.pptx"
OUT=DL/"NAMMA_KSP_Premium_Submission_Deck.pptx"
ARCH=Path(r"C:\Users\abdul\.codex\generated_images\019ed7ee-7b38-79d2-9f2b-ee25238bf127\exec-8f3c3741-c3e6-4474-be82-0cca54274c96.png")
SHOT=ROOT/"docs"/"screenshots"
NAVY=RGBColor(16,42,86); BLUE=RGBColor(31,90,166); TEAL=RGBColor(11,124,134)
GOLD=RGBColor(214,167,44); RED=RGBColor(190,31,40); INK=RGBColor(28,42,62)
MUTED=RGBColor(93,108,128); PALE=RGBColor(241,245,249); WHITE=RGBColor(255,255,255)
prs=Presentation(SRC)

def rm(shape): shape._element.getparent().remove(shape._element)
def clear(slide):
    for sh in list(slide.shapes)[1:]: rm(sh)
def rect(s,x,y,w,h,fill=WHITE,line=RGBColor(212,223,234),radius=True):
    q=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    q.fill.solid(); q.fill.fore_color.rgb=fill; q.line.color.rgb=line; q.line.width=Pt(.75); return q
def txt(s,v,x,y,w,h,size=14,color=INK,bold=False,align=PP_ALIGN.LEFT):
    q=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); f=q.text_frame; f.clear(); f.word_wrap=True
    f.margin_left=f.margin_right=Inches(.02); f.margin_top=f.margin_bottom=0; f.vertical_anchor=MSO_ANCHOR.TOP
    p=f.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=v; r.font.name="Aptos Display" if bold else "Aptos"; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return q
def pill(s,v,x,y,w,fill=GOLD,color=NAVY):
    q=rect(s,x,y,w,.29,fill,fill); q.text_frame.clear(); q.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=q.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=v; r.font.name="Aptos"; r.font.size=Pt(8); r.font.bold=True; r.font.color.rgb=color
def head(s,k,h,sub=""):
    pill(s,k.upper(),.62,.83,1.2,GOLD,NAVY)
    txt(s,h,.62,1.23,8.72,.55,21 if len(h)<58 else 18,NAVY,True)
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(.62),Inches(1.84),Inches(.72),Inches(.035)).fill.solid()
    line=s.shapes[-1]; line.fill.fore_color.rgb=TEAL; line.line.fill.background()
    if sub: txt(s,sub,1.48,1.75,7.85,.25,9,MUTED)
def pic(s,path,x,y,w,h):
    im=Image.open(path); iw,ih=im.size; ratio=iw/ih; target=w/h
    q=s.shapes.add_picture(str(path),Inches(x),Inches(y),width=Inches(w),height=Inches(h))
    if ratio>target: q.crop_left=q.crop_right=(1-target/ratio)/2
    else: q.crop_top=q.crop_bottom=(1-ratio/target)/2
    return q
def footer(s,n):
    txt(s,f"NAMMA KSP  /  KSP DATATHON 2026",.65,5.32,3,.12,6,MUTED,True)
    txt(s,f"{n:02}",9.05,5.27,.3,.16,7,NAVY,True,PP_ALIGN.RIGHT)

# Global headline corrections and premium editorial accents.
short={
2:"Natural-language questions. Evidence-backed action.",
3:"An intelligence workflow, not another dashboard",
4:"Ten challenge pillars. One operational platform.",
5:"Question → evidence → decision → accountable action",
8:"A pragmatic stack for speed, evidence and scale",
9:"Catalyst services mapped to operational value",
10:"Prototype economics with a clean scale path",
12:"Measured reliability, not demo theatre",
13:"Code, deployment and demonstration",
14:"From prototype to governed statewide intelligence",
15:"Prototype evidence and production guardrails",
}
for i in range(1,15):
    s=prs.slides[i]
    # Find the largest dark headline created by the first build.
    candidates=[sh for sh in s.shapes if hasattr(sh,"text_frame") and sh.text.strip() and sh.top>Inches(1.0) and sh.top<Inches(2.1)]
    if candidates:
        h=max(candidates,key=lambda z:z.width*z.height)
        if i+1 in short: h.text_frame.paragraphs[0].runs[0].text=short[i+1]
        h.left=Inches(.62); h.top=Inches(1.23); h.width=Inches(8.72); h.height=Inches(.53)
        for p in h.text_frame.paragraphs:
            for r in p.runs: r.font.size=Pt(20 if len(h.text)<58 else 18); r.font.bold=True; r.font.color.rgb=NAVY
    # Remove overlapping subtitle lines in the title zone; visual slides need air.
    for sh in list(s.shapes):
        if hasattr(sh,"text_frame") and sh is not (candidates and max(candidates,key=lambda z:z.width*z.height)) and sh.top>Inches(1.72) and sh.top<Inches(2.12) and sh.height<Inches(.5):
            if len(sh.text.strip())>40: rm(sh)
    # Stable footer and slim visual spine.
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(.28),Inches(.84),Inches(.035),Inches(4.37)); bar.fill.solid(); bar.fill.fore_color.rgb=TEAL; bar.line.fill.background()
    footer(s,i+1)

# Slide 2: editorial hero instead of four competing metric cards.
s=prs.slides[1]; clear(s); head(s,"Solution","Natural-language questions. Evidence-backed action.","One governed workspace for investigators, analysts, supervisors and policymakers.")
rect(s,.62,2.12,5.35,2.75,NAVY,NAVY); pill(s,"THE PROPOSITION",.88,2.4,1.25,GOLD,NAVY)
txt(s,"Turn fragmented crime records into connected, explainable intelligence.",.88,2.9,4.7,.86,25,WHITE,True)
txt(s,"Ask in English or Kannada. Trace people, places and incidents. Surface risk, evidence and next actions.",.9,4.0,4.55,.55,12,RGBColor(221,232,244))
for i,(n,l,c) in enumerate([("5,000","FIRs",BLUE),("7","linked datasets",TEAL),("2","languages + voice",GOLD)]):
    x=6.25; y=2.12+i*.91; rect(s,x,y,3.08,.73,WHITE,RGBColor(205,218,232)); txt(s,n,x+.18,y+.1,.9,.3,21,c,True); txt(s,l,x+1.08,y+.18,1.72,.25,11,NAVY,True)
rect(s,6.25,4.85,3.08,.03,RED,RED,False); footer(s,2)

# Slide 6: two large screenshots, no tiny triptych.
s=prs.slides[5]; clear(s); head(s,"Experience","Operational intelligence, designed for real investigative work")
for path,x,label,c in [(SHOT/'02-dashboard-strategic-briefing.png',.62,"STRATEGIC DASHBOARD",BLUE),(SHOT/'03-ai-chat.png',5.05,"BILINGUAL CONVERSATIONAL AI",TEAL)]:
    rect(s,x,2.08,4.32,2.78,WHITE,RGBColor(198,215,230)); pic(s,path,x+.08,2.16,4.16,2.2); pill(s,label,x+.25,4.5,3.82,c,WHITE)
footer(s,6)

# Slide 7: architecture becomes the hero and stays readable.
s=prs.slides[6]; clear(s); head(s,"Architecture","Governed intelligence from interaction to early warning")
pic(s,ARCH,.42,1.95,9.16,3.28); footer(s,7)

# Slide 11: four large real screenshots, edge-to-edge within disciplined grid.
s=prs.slides[10]; clear(s); head(s,"Prototype","Four working views. One connected investigation.")
items=[('04-crime-heatmap.png',.62,2.02,"HOTSPOT INTELLIGENCE",BLUE),('05-network-analysis.png',5.05,2.02,"CRIMINAL NETWORK",TEAL),('06-offender-profiles.png',.62,3.68,"OFFENDER INTELLIGENCE",GOLD),('07-reports.png',5.05,3.68,"REPORTS + EVIDENCE",RED)]
for fn,x,y,l,c in items:
    rect(s,x,y,4.32,1.48,WHITE,RGBColor(204,217,231)); pic(s,SHOT/fn,x+.06,y+.06,3.28,1.36); pill(s,l,x+3.43,y+.25,.76,c,WHITE)
footer(s,11)

# Slide 15: replace colliding title with a strong split-screen integrity statement.
s=prs.slides[14]; clear(s); head(s,"Integrity","Prototype evidence and production guardrails")
rect(s,.62,2.08,4.2,2.86,RGBColor(238,248,244),RGBColor(190,224,209)); rect(s,5.06,2.08,4.3,2.86,RGBColor(253,241,242),RGBColor(235,197,201))
pill(s,"VERIFIED NOW",.88,2.34,1.2,TEAL,WHITE); pill(s,"CONTROLLED NEXT",5.32,2.34,1.35,RED,WHITE)
left=["Working deployed web app","Synthetic linked datasets","Published QuickML endpoint","Groq + Sarvam integrations","PDF/QR and audit workflow"]
right=["No real KSP production data claim","Native auth configuration","Governed data onboarding","Human review of model outputs","Security and legal assessment"]
for i,v in enumerate(left): pill(s,"✓",.9,2.9+i*.36,.3,TEAL,WHITE); txt(s,v,1.34,2.91+i*.36,3.1,.22,10,INK,True)
for i,v in enumerate(right): pill(s,"!",5.34,2.9+i*.36,.3,RED,WHITE); txt(s,v,5.78,2.91+i*.36,3.18,.22,10,INK,True)
footer(s,15)

prs.save(OUT)
print(OUT)
